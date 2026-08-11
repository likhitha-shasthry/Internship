import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation(json_file_path, output_file_path):
    # Load JSON data
    with open(json_file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    course = data['course']
    chapter = course['chapter']

    # Create presentation object
    prs = Presentation()

    # Define common layout indices
    TITLE_SLIDE_LAYOUT = 0
    TITLE_AND_CONTENT_LAYOUT = 1
    SECTION_HEADER_LAYOUT = 2
    TWO_CONTENT_LAYOUT = 3

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Chapter {chapter['chapterNumber']}: {chapter['title']}"
    subtitle.text = f"{course['name']} | Grade {course['grade']} | Level: {course['level']}\n\n{chapter['description']}"
    
    # Optional styling for title
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].font.bold = True

    # --- Agenda/Navigation Slide ---
    slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Chapter Outline"
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for section in chapter['sections']:
        p = tf.add_paragraph()
        p.text = f"{section['number']} {section['title']}"
        p.level = 0

    # --- Sections ---
    for section in chapter['sections']:
        # Section Header
        slide_layout = prs.slide_layouts[SECTION_HEADER_LAYOUT]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{section['number']} {section['title']}"

        # Topics
        if 'topics' in section:
            for topic in section['topics']:
                slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = topic['title']
                
                tf = slide.placeholders[1].text_frame
                tf.clear()
                
                if 'content' in topic:
                    p = tf.add_paragraph()
                    p.text = topic['content']
                if 'definition' in topic:
                    p = tf.add_paragraph()
                    p.text = f"Definition: {topic['definition']}"
                if 'formula' in topic:
                    p = tf.add_paragraph()
                    p.text = f"Formula: {topic['formula']}"
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(153, 0, 0)
                if 'examples' in topic:
                    for ex in topic['examples']:
                        p = tf.add_paragraph()
                        p.text = f"• {ex}"
                        p.level = 1
                if 'condition' in topic and 'result' in topic:
                    p = tf.add_paragraph()
                    p.text = f"Condition: {topic['condition']}"
                    p = tf.add_paragraph()
                    p.text = f"Result: {topic['result']}"

        # Subsections
        if 'subsections' in section:
            for sub in section['subsections']:
                slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = sub['title']
                
                tf = slide.placeholders[1].text_frame
                tf.clear()
                
                if 'formula' in sub:
                    p = tf.add_paragraph()
                    p.text = f"Formula: {sub['formula']}"
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(153, 0, 0)
                if 'importantPoint' in sub:
                    p = tf.add_paragraph()
                    p.text = f"Note: {sub['importantPoint']}"
                if 'variables' in sub:
                    p = tf.add_paragraph()
                    p.text = "Variables:"
                    for k, v in sub['variables'].items():
                        p2 = tf.add_paragraph()
                        p2.text = f"{k}: {v}"
                        p2.level = 1

        # Example
        if 'example' in section:
            ex = section['example']
            slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"Example {ex['number']}: {ex['title']}"
            
            tf = slide.placeholders[1].text_frame
            tf.clear()
            
            if 'given' in ex:
                if isinstance(ex['given'], str):
                     p = tf.add_paragraph()
                     p.text = f"Given: {ex['given']}"
                else:
                     p = tf.add_paragraph()
                     p.text = "Given:"
                     for k, v in ex['given'].items():
                         p2 = tf.add_paragraph()
                         p2.text = f"{k}: {v}"
                         p2.level = 1
            if 'find' in ex:
                p = tf.add_paragraph()
                p.text = f"Find: {', '.join(ex['find'])}"
            if 'result' in ex:
                p = tf.add_paragraph()
                p.text = f"Result: {ex['result']}"
            if 'results' in ex:
                p = tf.add_paragraph()
                p.text = "Results:"
                for k, v in ex['results'].items():
                     p2 = tf.add_paragraph()
                     p2.text = f"{k}: {v}"
                     p2.level = 1
            if 'questions' in ex:
                for idx, q in enumerate(ex['questions']):
                     p = tf.add_paragraph()
                     p.text = f"Q{idx+1}: {q['question']}"
                     p2 = tf.add_paragraph()
                     p2.text = f"A{idx+1}: {q['answer']}"
                     p2.level = 1

    # --- Summary ---
    if 'summary' in chapter:
        slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Summary"
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        # PPT text frames can only hold so much text, so limit summary or split slides
        # We will split if > 10 items
        items = chapter['summary']
        for item in items[:10]:
            p = tf.add_paragraph()
            p.text = item
            
        if len(items) > 10:
             slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
             slide = prs.slides.add_slide(slide_layout)
             title = slide.shapes.title
             title.text = "Summary (Cont.)"
             
             tf = slide.placeholders[1].text_frame
             tf.clear()
             for item in items[10:]:
                 p = tf.add_paragraph()
                 p.text = item

    prs.save(output_file_path)
    print(f"Successfully generated PowerPoint: {output_file_path}")

if __name__ == "__main__":
    create_presentation('writing-block.json', 'Waves_Chapter_14.pptx')
